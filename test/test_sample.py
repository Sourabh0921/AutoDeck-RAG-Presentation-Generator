from pptx import Presentation


class PPTLoader:

    def extract_text(self, ppt_path):

        presentation = Presentation(ppt_path)

        slides_data = []

        for slide_number, slide in enumerate(presentation.slides, start=1):

            slide_text = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        slide_text.append(text)

            slides_data.append({
                "slide_number": slide_number,
                "content": "\n".join(slide_text)
            })

        return slides_data